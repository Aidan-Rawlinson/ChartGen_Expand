"""
template_reader.py
Reads a .pptx template and produces chart placeholders, extracted yellow-box content, and a cleaned template copy.
"""

import re
import os
from colorsys import rgb_to_hsv
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Hue range for "human yellow" in degrees
YELLOW_HUE_MIN = 40
YELLOW_HUE_MAX = 70
YELLOW_SAT_MIN = 0.50   # saturation floor — rejects cream / off-white
YELLOW_VAL_MIN = 0.50   # value floor — rejects dark mustard

# Toolkit URL pattern — matches NHS Benchmarking members URLs
TOOLKIT_URL_RE = re.compile(
    r"https?://[^\s\"'<>]+"
    r"nhsbenchmarking[^\s\"'<>]*",
    re.IGNORECASE,
)

# Fallback: any http/https URL
GENERIC_URL_RE = re.compile(r"https?://[^\s\"'<>]+", re.IGNORECASE)

# Image file extensions recognised by insert_picture
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".wmf", ".emf"}

# Excel workbook pattern: full path ending in .xlsx or .xlsm,
# followed by [key:rangename,...] bracket block
# e.g. C:\Projects\analysis.xlsx[driver:UNIT,export:SummaryTable]
EXCEL_PATH_RE = re.compile(
    r"^(.+\.(?:xlsx|xlsm))\[([^\]]+)\]$",
    re.IGNORECASE,
)

# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------

@dataclass
class PlaceholderInfo:
    """Describes a single named placeholder on a slide, OR a free-floating
    yellow box being used in its place (see read_template Step 3)."""
    slide_index: int          # 0-based
    name: str                 # e.g. "Chart 1" (placeholder) or the yellow box's own shape name (free-floating)
    left: int                 # EMU
    top: int                  # EMU
    width: int                # EMU
    height: int                # EMU
    # Yellow box classification — one of: "chart", "picture", "excel", ""
    content_type: str = ""
    # chart: toolkit URL
    url: str = ""
    label: str = ""           # user text from the yellow textbox (Running Order notes)
    # picture: image file path (may contain [code] or [id] tokens)
    image_path: str = ""
    # excel: workbook path, export range name, driver range name (optional)
    excel_path: str = ""
    excel_export_range: str = ""
    excel_driver_range: str = ""


@dataclass
class TemplateReadResult:
    """Full output of a template read operation."""
    slide_width: int          # EMU
    slide_height: int         # EMU
    placeholders: list = field(default_factory=list)   # list[PlaceholderInfo]
    warnings: list = field(default_factory=list)       # list[str]
    cleaned_pptx_bytes: bytes = b""                    # template with yellow boxes removed


# ---------------------------------------------------------------------------
# Colour helpers
# ---------------------------------------------------------------------------

def _rgb_to_hsv_degrees(r, g, b):
    """Convert 0–255 RGB to (hue_degrees, sat_0_1, val_0_1)."""
    h, s, v = rgb_to_hsv(r / 255, g / 255, b / 255)
    return h * 360, s, v


def _is_yellow(rgb) -> bool:
    """Return True if the RGB colour is a human-visible yellow."""
    if rgb is None:
        return False
    try:
        r, g, b = rgb[0], rgb[1], rgb[2]
    except (TypeError, IndexError):
        return False
    h, s, v = _rgb_to_hsv_degrees(r, g, b)
    return (
        YELLOW_HUE_MIN <= h <= YELLOW_HUE_MAX
        and s >= YELLOW_SAT_MIN
        and v >= YELLOW_VAL_MIN
    )


def _get_shape_fill_rgb(shape):
    """Return the effective solid fill RGBColor of a shape, or None if not
    determinable.

    Checked in order:
      1. An explicit <a:noFill/> on the shape itself — no fill, full stop;
         the shape's style reference (below) is not consulted.
      2. An explicit solid fill on the shape itself — literal RGB
         (<a:srgbClr>) or a theme colour reference (<a:schemeClr>).
      3. If the shape defines no fill of its own at all: its style's fill
         reference (<p:style><a:fillRef><a:schemeClr>) — the "Shape Styles"
         gallery mechanism, where the shape stores no colour, only a pointer
         to a theme colour slot. See Architecture, Decision 14.
      4. Fallback: python-pptx's native fill accessor, for anything the
         above didn't catch.
    """
    try:
        xml = shape._element.xml
    except Exception:
        return None

    spPr_m = re.search(r"<p:spPr>.*?</p:spPr>", xml, re.DOTALL)
    spPr_xml = spPr_m.group(0) if spPr_m else ""

    # --- 1. Explicit no-fill on the shape itself ---
    if re.search(r"<a:noFill\s*/>", spPr_xml):
        return None

    # --- 2a. Explicit literal RGB fill on the shape itself ---
    m = re.search(r'<a:solidFill>\s*<a:srgbClr\s+val="([0-9A-Fa-f]{6})"', spPr_xml)
    if m:
        hex_val = m.group(1).upper()
        from pptx.dml.color import RGBColor
        return RGBColor(int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16))

    # --- 2b. Explicit theme-colour fill on the shape itself ---
    m = re.search(r'<a:solidFill>\s*<a:schemeClr\s+val="(\w+)"', spPr_xml)
    if m:
        clr_map, clr_scheme = _get_theme_context(shape)
        rgb = _resolve_scheme_colour(m.group(1), clr_map, clr_scheme)
        if rgb:
            from pptx.dml.color import RGBColor
            return RGBColor(*rgb)

    # --- 3. No fill defined on the shape at all: fall back to its style's fillRef ---
    if "<a:solidFill" not in spPr_xml:
        style_m = re.search(r"<p:style>.*?</p:style>", xml, re.DOTALL)
        if style_m:
            fillref_m = re.search(r'<a:fillRef[^>]*>\s*<a:schemeClr\s+val="(\w+)"', style_m.group(0))
            if fillref_m:
                clr_map, clr_scheme = _get_theme_context(shape)
                rgb = _resolve_scheme_colour(fillref_m.group(1), clr_map, clr_scheme)
                if rgb:
                    from pptx.dml.color import RGBColor
                    return RGBColor(*rgb)

    # --- 4. Fallback: python-pptx's native accessor ---
    try:
        fill = shape.fill
        if str(fill.type) == "SOLID (1)":
            return fill.fore_color.rgb
    except Exception:
        pass
    return None


# ---------------------------------------------------------------------------
# Theme colour resolution
# ---------------------------------------------------------------------------

# The 12 base colour slots every theme's <a:clrScheme> defines.
_THEME_SLOTS = ("dk1", "lt1", "dk2", "lt2",
                "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
                "hlink", "folHlink")

# Slot names a slide's/master's <p:clrMap> (or a slide's <p:clrMapOvr>) can remap.
_CLR_MAP_ATTRS = ("bg1", "tx1", "bg2", "tx2",
                   "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
                   "hlink", "folHlink")

# Cache: slide master id -> (clr_map: dict, clr_scheme: dict), avoids re-parsing
# the master's colour map and its theme's colour scheme for every shape.
_theme_context_cache: dict = {}


def _parse_clr_scheme(theme_xml_bytes: bytes) -> dict:
    """Parse a theme part's <a:clrScheme> into {slot_name: (r, g, b)}."""
    scheme = {}
    text = theme_xml_bytes.decode("utf-8", errors="ignore")
    m = re.search(r"<a:clrScheme[^>]*>(.*?)</a:clrScheme>", text, re.DOTALL)
    if not m:
        return scheme
    body = m.group(1)
    for slot in _THEME_SLOTS:
        sm = re.search(rf'<a:{slot}>\s*<a:srgbClr val="([0-9A-Fa-f]{{6}})"', body)
        if not sm:
            # A theme colour can also be tied to a system colour (rare, e.g. window/windowText)
            sm = re.search(rf'<a:{slot}>\s*<a:sysClr[^>]*lastClr="([0-9A-Fa-f]{{6}})"', body)
        if sm:
            hex_val = sm.group(1)
            scheme[slot] = (int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16))
    return scheme


def _parse_clr_map(clr_map_element_xml: str) -> dict:
    """Parse a <p:clrMap .../> or <a:overrideClrMapping .../> element's
    attributes into {attribute_name: resolved_theme_slot_name}."""
    clr_map = {}
    for attr in _CLR_MAP_ATTRS:
        m = re.search(rf'\b{attr}="(\w+)"', clr_map_element_xml)
        if m:
            clr_map[attr] = m.group(1)
    return clr_map


def _get_theme_context(shape):
    """Return (clr_map, clr_scheme) for the slide a shape sits on: the
    colour-name remapping in effect (slide-level override if present, else
    the slide master's own map) and the theme colour scheme belonging to
    that master. Cached per master to avoid repeated XML parsing."""
    try:
        slide = shape.part.slide
        master = slide.slide_layout.slide_master
    except Exception:
        return {}, {}

    cache_key = id(master)
    if cache_key not in _theme_context_cache:
        master_xml = master._element.xml

        map_m = re.search(r"<p:clrMap\b[^/]*/>", master_xml)
        base_clr_map = _parse_clr_map(map_m.group(0)) if map_m else {}

        clr_scheme = {}
        try:
            theme_part = master.part.part_related_by(RELATIONSHIP_TYPE.THEME)
            clr_scheme = _parse_clr_scheme(theme_part.blob)
        except Exception:
            pass

        _theme_context_cache[cache_key] = (base_clr_map, clr_scheme)

    base_clr_map, clr_scheme = _theme_context_cache[cache_key]

    # A slide can override its master's colour map wholesale via <p:clrMapOvr>.
    try:
        slide_xml = slide._element.xml
    except Exception:
        slide_xml = ""
    ovr_m = re.search(r"<p:clrMapOvr>.*?</p:clrMapOvr>", slide_xml, re.DOTALL)
    if ovr_m and "overrideClrMapping" in ovr_m.group(0):
        override_m = re.search(r"<a:overrideClrMapping\b[^/]*/>", ovr_m.group(0))
        clr_map = _parse_clr_map(override_m.group(0)) if override_m else base_clr_map
    else:
        clr_map = base_clr_map

    return clr_map, clr_scheme


def _resolve_scheme_colour(scheme_name: str, clr_map: dict, clr_scheme: dict):
    """Resolve a <a:schemeClr val="..."/> name to an (r, g, b) tuple, or None
    if it can't be resolved (e.g. "phClr" — a contextual placeholder colour
    with no fixed value outside the effect/style definition referencing it)."""
    if scheme_name == "phClr":
        return None
    resolved_slot = clr_map.get(scheme_name, scheme_name)
    return clr_scheme.get(resolved_slot)


def _extract_url_from_text(text: str) -> str:
    """Extract the first toolkit URL from text, falling back to any HTTP URL; returns empty string if none found."""
    m = TOOLKIT_URL_RE.search(text)
    if m:
        return m.group(0).rstrip(".,;)")
    m = GENERIC_URL_RE.search(text)
    if m:
        return m.group(0).rstrip(".,;)")
    return ""


def _classify_yellow_box(text: str) -> dict:
    """
    Classify the content of a yellow textbox and return a dict describing it.

    Returns a dict with key 'content_type' set to one of:
      "chart"   — text contains a toolkit or HTTP URL
      "picture" — text is a path to a supported image file
      "excel"   — text matches the Excel path + range bracket syntax
      ""        — unrecognised; yellow box will be stripped but ignored

    Additional keys depend on content_type:
      chart:   url, label
      picture: image_path
      excel:   excel_path, excel_export_range, excel_driver_range
    """
    text = text.strip()

    # --- Excel: path.xlsx[driver:X,export:Y] ---
    m = EXCEL_PATH_RE.match(text)
    if m:
        excel_path  = m.group(1).strip()
        params_str  = m.group(2).strip()
        params = {}
        for part in params_str.split(","):
            part = part.strip()
            if ":" in part:
                key, val = part.split(":", 1)
                params[key.strip().lower()] = val.strip()
        export_range = params.get("export", "")
        driver_range = params.get("driver", "")
        if export_range:
            return {
                "content_type":       "excel",
                "excel_path":         excel_path,
                "excel_export_range": export_range,
                "excel_driver_range": driver_range,
            }

    # --- Picture: path ending in a supported image extension ---
    # Strip [code] and [id] tokens before checking extension
    clean = text.replace("[code]", "X").replace("[id]", "X")
    ext = os.path.splitext(clean)[1].lower()
    if ext in IMAGE_EXTENSIONS:
        return {
            "content_type": "picture",
            "image_path":   text,
        }

    # --- Chart: URL ---
    url = _extract_url_from_text(text)
    if url:
        return {
            "content_type": "chart",
            "url":          url,
            "label":        text.strip(),
        }

    return {"content_type": ""}


# ---------------------------------------------------------------------------
# Geometry helpers
# ---------------------------------------------------------------------------

# 1mm of tolerance on the "fully contained" check, to absorb sub-visible EMU
# rounding drift (e.g. PowerPoint copy/paste introducing a 1 EMU discrepancy
# on a duplicated shape) without treating it as a genuine partial overlap.
# 914400 EMU = 1 inch = 25.4mm, so 36000 EMU = 1mm exactly.
CONTAINMENT_TOLERANCE_EMU = 36000

def _fully_contained(inner_left, inner_top, inner_right, inner_bottom,
                     outer_left, outer_top, outer_right, outer_bottom,
                     tolerance=CONTAINMENT_TOLERANCE_EMU) -> bool:
    """Return True if inner rectangle is within outer rectangle, allowing
    up to `tolerance` EMU of drift on each edge."""
    return (
        inner_left   >= outer_left   - tolerance
        and inner_top    >= outer_top    - tolerance
        and inner_right  <= outer_right  + tolerance
        and inner_bottom <= outer_bottom + tolerance
    )


def _rects_overlap(a_left, a_top, a_right, a_bottom,
                    b_left, b_top, b_right, b_bottom) -> bool:
    """Return True if two rectangles share any non-zero area. Touching edges
    only (no interior overlap) do not count."""
    return (
        a_left < b_right and a_right > b_left
        and a_top < b_bottom and a_bottom > b_top
    )


# ---------------------------------------------------------------------------
# Shape classification
# ---------------------------------------------------------------------------

def _is_chart_placeholder(shape) -> bool:
    """Return True if this shape is a placeholder ChartGen should treat as a chart slot."""
    if not shape.is_placeholder:
        return False
    ph_type = shape.placeholder_format.type
    _content_types = {
        PP_PLACEHOLDER.OBJECT,
        PP_PLACEHOLDER.PICTURE,
        PP_PLACEHOLDER.CHART,
        PP_PLACEHOLDER.BITMAP,
        PP_PLACEHOLDER.TABLE,
        PP_PLACEHOLDER.ORG_CHART,
        PP_PLACEHOLDER.MEDIA_CLIP,
    }
    return ph_type in _content_types


def _is_textbox(shape) -> bool:
    """Return True if shape is a free textbox (not a placeholder)."""
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    except Exception:
        return False


def _is_autoshape_with_text(shape) -> bool:
    """Return True if shape is an autoshape (rectangle etc.) with text — also counts."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE as MST
        return shape.shape_type == MST.AUTO_SHAPE and shape.has_text_frame
    except Exception:
        return False


# ---------------------------------------------------------------------------
# Core read function
# ---------------------------------------------------------------------------

def read_template(pptx_path: str) -> TemplateReadResult:
    """Read a .pptx template and return a TemplateReadResult.

    A yellow textbox is resolved against the slide's chart placeholders into
    one of three outcomes (see Architecture, Decision 13):

      1. Fully contained by a placeholder  — matched to it; the placeholder's
         own position/size is used (the pre-existing behaviour).
      2. No overlap with any placeholder   — free-floating; the yellow box's
         own position/size is used directly, named after its own shape name.
      3. Partial overlap with a placeholder, short of full containment —
         ambiguous; left entirely untouched (not processed, not removed).
    """
    prs = Presentation(pptx_path)
    result = TemplateReadResult(
        slide_width=int(prs.slide_width),
        slide_height=int(prs.slide_height),
    )

    # We need to track which XML elements to remove for the cleaned copy.
    # Work on the live prs first, collect removal targets, then strip.
    elements_to_remove = []   # list of (slide, sp_element)

    for slide_idx, slide in enumerate(prs.slides):
        # --- Step 1: collect chart placeholders on this slide ---
        chart_placeholders = []
        ph_shapes: dict[str, object] = {}   # name -> shape, for removal once matched
        for shape in slide.shapes:
            if _is_chart_placeholder(shape):
                chart_placeholders.append(
                    PlaceholderInfo(
                        slide_index=slide_idx,
                        name=shape.name,
                        left=int(shape.left),
                        top=int(shape.top),
                        width=int(shape.width),
                        height=int(shape.height),
                    )
                )
                ph_shapes[shape.name] = shape

        # --- Step 2: collect yellow textboxes on this slide ---
        yellow_shapes = []
        for shape in slide.shapes:
            if not (_is_textbox(shape) or _is_autoshape_with_text(shape)):
                continue
            rgb = _get_shape_fill_rgb(shape)
            if not _is_yellow(rgb):
                continue
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            classification = _classify_yellow_box(text)
            if not classification.get("content_type"):
                # Unrecognised content — doesn't match any of the four
                # supported types (NHS toolkit chart URL, Indicators toolkit
                # chart URL, image path, Excel path+ranges). Stripped from
                # the cleaned template as before (Functional Spec §6.3), but
                # now flagged rather than silently dropped.
                preview = text.strip()
                if len(preview) > 80:
                    preview = preview[:77] + "..."
                result.warnings.append(
                    f"Slide {slide_idx + 1}: yellow textbox content didn't match "
                    f"a recognised type (NHS chart URL, Indicators chart URL, "
                    f"picture path, or Excel path+ranges) — stripped, not processed. "
                    f"Text: \"{preview}\""
                )
                elements_to_remove.append((slide, shape._element))
                continue
            yellow_shapes.append({
                "shape":  shape,
                "left":   int(shape.left),
                "top":    int(shape.top),
                "right":  int(shape.left) + int(shape.width),
                "bottom": int(shape.top)  + int(shape.height),
                **classification,
            })

        # --- Step 3: resolve each yellow textbox against the placeholders ---
        ph_matches: dict[str, list] = {ph.name: [] for ph in chart_placeholders}
        free_floating = []   # list[PlaceholderInfo] — scenario 2, one per free-floating box

        for yshape in yellow_shapes:
            contained_by = None
            overlaps_any = False
            for ph in chart_placeholders:
                ph_right  = ph.left + ph.width
                ph_bottom = ph.top  + ph.height
                if _fully_contained(
                    yshape["left"], yshape["top"],
                    yshape["right"], yshape["bottom"],
                    ph.left, ph.top, ph_right, ph_bottom,
                ):
                    contained_by = ph
                    break
                if _rects_overlap(
                    yshape["left"], yshape["top"],
                    yshape["right"], yshape["bottom"],
                    ph.left, ph.top, ph_right, ph_bottom,
                ):
                    overlaps_any = True

            if contained_by is not None:
                # Scenario 1 — fully contained.
                ph_matches[contained_by.name].append(yshape)
            elif overlaps_any:
                # Scenario 3 — ambiguous partial overlap: not processed, not removed.
                yshape["_skip"] = True
                result.warnings.append(
                    f"Slide {slide_idx + 1}: yellow textbox partially overlaps a "
                    f"placeholder without being fully inside it — not processed. "
                    f"Content: {yshape.get('url') or yshape.get('image_path') or yshape.get('excel_path', '')}"
                )
            else:
                # Scenario 2 — free-floating: use the yellow box's own position/size.
                free_floating.append(
                    PlaceholderInfo(
                        slide_index=slide_idx,
                        name=yshape["shape"].name,
                        left=yshape["left"],
                        top=yshape["top"],
                        width=yshape["right"]  - yshape["left"],
                        height=yshape["bottom"] - yshape["top"],
                        content_type=yshape.get("content_type", ""),
                        url=yshape.get("url", ""),
                        label=yshape.get("label", ""),
                        image_path=yshape.get("image_path", ""),
                        excel_path=yshape.get("excel_path", ""),
                        excel_export_range=yshape.get("excel_export_range", ""),
                        excel_driver_range=yshape.get("excel_driver_range", ""),
                    )
                )

        # --- Step 4: assign content to matched placeholders (scenario 1) ---
        for ph in chart_placeholders:
            matches = ph_matches[ph.name]
            if len(matches) == 0:
                pass
            else:
                if len(matches) > 1:
                    result.warnings.append(
                        f"Slide {slide_idx + 1}, placeholder '{ph.name}': "
                        f"{len(matches)} yellow textboxes found — using first."
                    )
                m = matches[0]
                ct = m["content_type"]
                ph.content_type = ct
                if ct == "chart":
                    ph.url   = m.get("url",   "")
                    ph.label = m.get("label", "")
                elif ct == "picture":
                    ph.image_path = m.get("image_path", "")
                elif ct == "excel":
                    ph.excel_path         = m.get("excel_path",         "")
                    ph.excel_export_range = m.get("excel_export_range", "")
                    ph.excel_driver_range = m.get("excel_driver_range", "")

                # This placeholder is now fully captured (position/size in
                # PlaceholderInfo, content assignment above) — content is
                # inserted at generation time by coordinate, not via the
                # placeholder shape itself (Functional Spec §6.2), so the
                # placeholder shape is removed from the cleaned template
                # alongside its yellow box, rather than left behind empty.
                ph_shape = ph_shapes.get(ph.name)
                if ph_shape is not None:
                    elements_to_remove.append((slide, ph_shape._element))

            result.placeholders.append(ph)

        # --- Step 4b: scenario 2 — free-floating boxes stand in as their own placeholders ---
        result.placeholders.extend(free_floating)

        # --- Step 5: mark yellow textboxes for removal (scenarios 1 & 2 only) ---
        for yshape in yellow_shapes:
            if yshape.get("_skip"):
                continue   # scenario 3 — ambiguous overlap, left in place untouched
            elements_to_remove.append((slide, yshape["shape"]._element))

    # --- Step 6: produce cleaned pptx bytes ---
    import io as _io
    for slide, elem in elements_to_remove:
        try:
            elem.getparent().remove(elem)
        except Exception:
            pass  # already removed or detached

    buf = _io.BytesIO()
    prs.save(buf)
    result.cleaned_pptx_bytes = buf.getvalue()

    if result.warnings:
        result.warnings.insert(
            0,
            "One or more yellow boxes could not be processed — see details below."
        )

    return result
