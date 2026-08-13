"""
assembly_engine.py
Executes a Running Order's normal-scope rows against an open PowerPoint
template to produce one report. Dispatches each row to its Running Order
function (create_ppt, insert_chart, insert_picture, insert_from_excel,
update_text, save_ppt, etc) via FUNCTION_MAP and returns a per-run log.

Note (Restructure_Plan.md Open Item 1): this module was previously described
as "the only package touching python-pptx directly" — that is no longer
true now insert_picture and insert_from_excel also manipulate python-pptx
objects. Its actual purpose is dispatch/execution of one report's Running
Order rows, not exclusive ownership of python-pptx.
"""

import os
import time
import traceback
from dataclasses import replace

from pptx import Presentation

from core.output_generation.execution.charts.cache_reader import load_shape
from core.output_generation.execution.charts.custom_charts import get_chart_callable
from core.output_generation.execution.text.text_engine import update_text
from core.output_generation.execution.tables.insert_table import insert_table
from core.output_generation.execution.svg_insert import add_svg_picture
from core.shared.infrastructure.report_context import build_report_context
from core.shared.infrastructure.soft_parents import resolve_full_unit_set
from core.shared.normalisation_containers.cut_resolution import prepare_chart_cut
from core.shared.normalisation_containers.population_layers import build_population_layers
from core.output_generation.execution.pictures.insert_picture import insert_picture
from core.output_generation.execution.excel.insert_from_excel import (
    open_excel, close_excel, insert_from_excel
)
from core.output_generation.execution.results import ok_result, err_result
from core.workfile.state.workfile_file import master_table_rows


# ---------------------------------------------------------------------------
# Hyperlink icon — optional, insert_chart only
# ---------------------------------------------------------------------------

DEFAULT_HYPERLINK_COLOUR = "#0563C1"   # standard Office hyperlink blue
DEFAULT_HYPERLINK_SIZE_EMU = 360000    # ~1cm (914400 EMU/inch ÷ 2.54)


def _hyperlink_icon_svg_bytes(colour_hex: str) -> bytes:
    """
    A small chain-link icon, drawn from scratch (two rounded-rect "links",
    not traced from any icon library) -- own copy, not shared with any
    Base Chart, since this is assembly-layer decoration, not a chart.
    The two links overlap (stroke-only rects, no fill, so the background
    shows through in the overlap band) rather than sitting apart with a
    gap -- matching how a real link symbol is read (the two rings joined,
    not separate). viewBox is a fixed 72x72 square regardless of the
    icon's actual placed size on the slide -- add_svg_picture scales it to
    whatever width_emu/height_emu it's given, same as every Base Chart's
    own SVG output.
    """
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 72 72">'
        f'<g fill="none" stroke="{colour_hex}" stroke-width="6" stroke-linecap="round">'
        '<rect x="14" y="26" width="26" height="20" rx="10"/>'
        '<rect x="32" y="26" width="26" height="20" rx="10"/>'
        '</g></svg>'
    )
    return svg.encode("utf-8")


def _insert_hyperlink_icon(prs: Presentation, slide_index: int,
                            chart_left_emu: int, chart_top_emu: int, chart_width_emu: int,
                            hyperlink_left_emu: int, hyperlink_top_emu: int,
                            size_emu: int, colour_hex: str, source_url: str = None):
    """
    Draw the hyperlink icon onto the slide, positioned relative to the
    chart's own top-right corner -- NOT an absolute slide position, so the
    icon travels with the chart regardless of where the chart itself sits.
    (hyperlink_left_emu, hyperlink_top_emu) of (0, 0) places the icon's own
    top-left corner exactly at the chart's top-right corner. Drawn square:
    size_emu is used for both width and height. Sub-step of insert_chart,
    called after the chart image itself is already on the slide.

    source_url, if given, is set as the icon shape's own click hyperlink
    (python-pptx's click_action.hyperlink.address) -- the icon's own
    identity as a link, not a hyperlink on the chart picture itself. A
    blank/missing source_url still draws the icon (position/colour are
    resolved independently of whether the data shape's own metadata has a
    URL recorded) -- it just isn't clickable.

    Returns the icon's own shape -- insert_chart names it (CG_Link_{row_id})
    once it's back in scope there, rather than this function knowing about
    Running Order row identity at all.
    """
    icon_left_emu = chart_left_emu + chart_width_emu + hyperlink_left_emu
    icon_top_emu = chart_top_emu + hyperlink_top_emu
    slide = prs.slides[slide_index]
    icon_shape = add_svg_picture(
        slide, _hyperlink_icon_svg_bytes(colour_hex),
        icon_left_emu, icon_top_emu, size_emu, size_emu,
    )
    if source_url:
        icon_shape.click_action.hyperlink.address = source_url
    return icon_shape


# ---------------------------------------------------------------------------
# Assembly context — passed through every function call in a run
# ---------------------------------------------------------------------------

class AssemblyContext:
    def __init__(self):
        self.prs: Presentation = None
        self.output_path: str = ""
        self.template_path: str = ""
        # No consumer reads this full list today — batch_process.py's
        # per-unit run log only surfaces the first error, now prefixed
        # with its own row_id (see results.py's err_result and
        # batch_process.py) so a failure is attributable to a specific
        # row even from a one-line summary. Left in deliberately rather
        # than stripped: a full per-row (function/status/message) log of
        # one report's run could be genuinely useful for future
        # debugging/diagnostics, if a real consumer is ever built.
        self.log: list[dict] = []
        self.report_context = None      # set by run_running_order
        self.full_unit_set: dict = {}   # {table_name: [row, ...]} for the current reporting unit, set by run_running_order
        self.default_populations: str = ""  # set by set_default_populations row


# ---------------------------------------------------------------------------
# Running Order function implementations
# ---------------------------------------------------------------------------

def create_ppt(ctx: AssemblyContext, row: dict, settings: dict) -> dict:
    """
    Open the cleaned template and set the output path.
    Does not write to disk — save_ppt/save_pdf do that.
    """
    template_path = settings.get("cleaned_template_path", "").strip()
    if not template_path or not os.path.exists(template_path):
        # Fall back to the original template if no cleaned version exists
        template_path = settings.get("ppt_template_path", "").strip()

    if not template_path or not os.path.exists(template_path):
        return err_result(row, "create_ppt: no template found. Check settings.")

    output_folder = _ensure_output_folder(settings)
    unit_name = (settings.get("reporting_unit_name") or "").strip()
    if not unit_name:
        unit_name = str(settings.get("selected_unit_id") or "output").strip()
    safe_name = _safe_filename(unit_name)
    output_path = os.path.join(output_folder, "pptx", f"{safe_name}.pptx")

    ctx.prs = Presentation(template_path)
    # Force autoCompressPictures="0" on the presentation's own root XML
    # element (ppt/presentation.xml) -- the same flag PowerPoint's "Do not
    # compress images in file" checkbox controls (ISO/IEC 29500-1,
    # section 19.2.1.26), stored per-presentation. python-pptx has no
    # dedicated property for this; set directly on the underlying element.
    # TEST, not yet a settled decision -- see Architecture, Structural
    # Design Principles ("Validate only where designed"): confirming
    # whether PowerPoint's own PDF export is silently downsampling images
    # before deciding whether this stays.
    ctx.prs.part._element.set("autoCompressPictures", "0")
    ctx.output_path = output_path
    ctx.template_path = template_path

    return ok_result(row, f"Template opened: {os.path.basename(template_path)}")


def set_default_populations(ctx: AssemblyContext, row: dict, settings: dict) -> dict:
    """
    Store the default populations string on AssemblyContext.
    Subsequent insert_chart rows inherit it unless overridden per row.
    """
    populations = str(row.get("populations", "") or "").strip()
    ctx.default_populations = populations
    return ok_result(row, f"Default populations set: '{populations}'")


def insert_chart(ctx: AssemblyContext, row: dict, settings: dict) -> dict:
    """Render a Base Chart from cached data and insert it at the row's position."""
    if ctx.prs is None:
        return err_result(row, "insert_chart: no open presentation (create_ppt not called?).")

    cache_file = str(row.get("cache_file") or "").strip()
    if cache_file.lower() == "none":
        cache_file = ""
    base_chart_name = str(row.get("base_chart_name", "")).strip()
    slide_index = _int_or_none(row.get("slide_index"))

    # Position / size from the Running Order row (written from template at generation time)
    left_emu = _int_or_none(row.get("left_emu"))
    top_emu = _int_or_none(row.get("top_emu"))
    width_emu = _int_or_none(row.get("width_emu"))
    height_emu = _int_or_none(row.get("height_emu"))

    # Validate required fields
    missing = []
    if not cache_file:  missing.append("cache_file")
    if not base_chart_name: missing.append("base_chart_name")
    if slide_index is None: missing.append("slide_index")
    if None in (left_emu, top_emu, width_emu, height_emu):
        missing.append("position/size EMU values")
    if missing:
        return err_result(row, f"insert_chart: missing required fields: {', '.join(missing)}")

    # --- Resolve populations for this chart ---
    row_populations = str(row.get("populations", "") or "").strip()
    populations_str = row_populations if row_populations else ctx.default_populations

    render_context = ctx.report_context

    # --- Load data shape ---
    try:
        data_shape, shape_type = _load_chart_data(cache_file, settings.get("workfile_state"))
    except Exception as e:
        return err_result(row, f"insert_chart: failed to load cache '{cache_file}': {e}")

    # --- Resolve this row's own cut of the data shape. Period-range trim,
    # metric-periods conversion, and population-table/target-rows/
    # selected-ids resolution are all composed in
    # cut_resolution.prepare_chart_cut, shared with the Charts sheet and
    # stat tags — see that module. data_shape comes back trimmed/converted
    # regardless of whether any layers are actually resolved below, so the
    # "no populations" fallback reflects those same trims rather than
    # reverting to the untrimmed shape. An unresolvable metric_periods id
    # doesn't raise (see time_series_to_numeric_series' own docstring) —
    # it comes through as a real metric with no data for any unit, for
    # the Base Chart itself to handle, the same as any other missing
    # value. ---
    start_period = str(row.get("start_period", "") or "").strip()
    end_period = str(row.get("end_period", "") or "").strip()
    metric_periods_str = str(row.get("metric_periods", "") or "").strip()
    workfile_state = settings.get("workfile_state")

    data_shape, _, target_rows, selected_ids = prepare_chart_cut(
        data_shape, shape_type, start_period, end_period, metric_periods_str,
        workfile_state.tables, workfile_state.table_order, ctx.full_unit_set,
    )

    population_layers = []
    if render_context is not None and populations_str:
        try:
            population_layers = build_population_layers(
                data_shape, populations_str, target_rows, selected_ids
            )
        except Exception as e:
            return err_result(row, f"insert_chart: failed to build population layers: {e}")

    # Fall back to full unfiltered shape if no populations resolved
    if not population_layers:
        population_layers = [replace(data_shape, population_label="All")]

    # --- Render chart image ---
    tweaks = str(row.get("tweaks", "") or "").strip()
    try:
        image_bytes = _render_chart_image(
            base_chart_name, population_layers, width_emu, height_emu, tweaks,
            settings.get("workfile_state").custom_chart_code,
        )
    except Exception as e:
        return err_result(row, f"insert_chart: render failed for '{base_chart_name}': {e}")

    # --- Insert into slide ---
    try:
        chart_shape = _insert_image_at_position(
            ctx.prs, slide_index,
            image_bytes, left_emu, top_emu, width_emu, height_emu
        )
        # Named for traceback against the Running Order (Position Finder
        # tool, running_order_tab.py) -- keyed on row_id, this row's
        # current line number. row_id is renumbered whenever rows are
        # inserted/reordered/deleted (row_ops.renumber_row_ids), so this
        # name only stays accurate until the Running Order is next
        # edited -- an accepted trade-off (Aidan's own call) rather than
        # a genuinely stable identity like hex_id or a Stat Tag id.
        chart_shape.name = f"CG_Chart_{row.get('row_id')}"
    except Exception as e:
        return err_result(row, f"insert_chart: failed to insert image on slide {slide_index}: {e}")

    # --- Optional hyperlink icon, positioned relative to the chart's own
    # top-right corner. Generates only when BOTH hyperlink_left and
    # hyperlink_top are present on the row -- blank in either means no
    # icon at all. (0, 0) is a valid, meaningful value distinct from
    # blank: it places the icon's own top-left corner exactly at the
    # chart's top-right corner. hyperlink_size/hyperlink_colour each fall
    # back to their own default independently of whether they're blank. ---
    hyperlink_left_raw = str(row.get("hyperlink_left", "") or "").strip()
    hyperlink_top_raw = str(row.get("hyperlink_top", "") or "").strip()
    if hyperlink_left_raw != "" and hyperlink_top_raw != "":
        hyperlink_left_emu = _int_or_none(hyperlink_left_raw)
        hyperlink_top_emu = _int_or_none(hyperlink_top_raw)
        if hyperlink_left_emu is None or hyperlink_top_emu is None:
            return err_result(row, "insert_chart: hyperlink_left/hyperlink_top must be whole EMU numbers.")
        hyperlink_size_emu = _int_or_none(row.get("hyperlink_size")) or DEFAULT_HYPERLINK_SIZE_EMU
        hyperlink_colour = str(row.get("hyperlink_colour", "") or "").strip() or DEFAULT_HYPERLINK_COLOUR
        # source_url comes from this row's own resolved data shape (the
        # cut actually rendered, data_shape, in scope above) -- not from
        # any particular population layer, since metadata isn't part of
        # any single layer's own filtered view; it's the same value on
        # every layer, carried through filter/replace() unchanged from
        # where fetch.py recorded it.
        source_url = (data_shape.metadata or {}).get("source_url")
        try:
            icon_shape = _insert_hyperlink_icon(
                ctx.prs, slide_index,
                left_emu, top_emu, width_emu,
                hyperlink_left_emu, hyperlink_top_emu,
                hyperlink_size_emu, hyperlink_colour,
                source_url=source_url,
            )
            # Same naming convention/caveat as the chart shape above.
            icon_shape.name = f"CG_Link_{row.get('row_id')}"
        except Exception as e:
            return err_result(row, f"insert_chart: failed to insert hyperlink icon on slide {slide_index}: {e}")

    return ok_result(row, f"Chart '{base_chart_name}' inserted (slide {slide_index + 1})")


def empty_placeholder(ctx: AssemblyContext, row: dict, settings: dict) -> dict:
    """No-op. Placeholder has no content assigned."""
    return ok_result(row, f"empty_placeholder: row {row.get('row_id')} skipped (no content assigned)")


def save_ppt(ctx: AssemblyContext, row: dict, settings: dict) -> dict:
    """Save the completed output as a .pptx file."""
    if ctx.prs is None:
        return err_result(row, "save_ppt: no open presentation.")
    try:
        os.makedirs(os.path.dirname(ctx.output_path), exist_ok=True)
        ctx.prs.save(ctx.output_path)
        return ok_result(row, f"Saved: {ctx.output_path}")
    except Exception as e:
        return err_result(row, f"save_ppt: {e}")


def save_pdf(ctx: AssemblyContext, row: dict, settings: dict) -> dict:
    """
    Save the completed output as a .pdf using COM automation (Windows/PowerPoint only).
    Falls back gracefully on non-Windows or if PowerPoint is not available.
    """
    if ctx.prs is None:
        return err_result(row, "save_pdf: no open presentation.")

    pdf_dir = os.path.join(os.path.dirname(os.path.dirname(ctx.output_path)), "pdf")
    pdf_path = os.path.join(pdf_dir, os.path.basename(ctx.output_path).replace(".pptx", ".pdf"))
    os.makedirs(pdf_dir, exist_ok=True)

    # Ensure the pptx is saved first (COM needs a file on disk to open)
    try:
        ctx.prs.save(ctx.output_path)
    except Exception as e:
        return err_result(row, f"save_pdf: could not save .pptx before PDF export: {e}")

    try:
        import pythoncom
        import comtypes.client
        # Explicit CoInitialize/CoUninitialize on this call, rather than
        # relying on comtypes' own implicit init -- comtypes tracks "have I
        # initialised COM?" as a module-level flag, not per-thread, and
        # Streamlit runs each script rerun on a fresh thread. A second run
        # landing on a new thread that's never actually had CoInitialize
        # called on it, but where comtypes' own flag already says "done",
        # skips initialisation and fails with "CoInitialize has not been
        # called" the moment the COM object is actually used. Same
        # explicit-call convention insert_from_excel.py already uses for
        # its own COM session, for the same reason.
        pythoncom.CoInitialize()
        try:
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            deck = powerpoint.Presentations.Open(os.path.abspath(ctx.output_path))
            # TEST, not yet a settled decision -- see Architecture, Structural
            # Design Principles ("Validate only where designed"). ExportAsFixedFormat
            # is the newer PDF export pathway (FixedFormatType=2 = ppFixedFormatTypePDF),
            # called with default settings for everything else. Decision 26 originally
            # moved away from this method because it produced visibly downsampled
            # embedded images even with autoCompressPictures forced off -- revisited
            # here at Aidan's request; that finding may still apply and is worth
            # re-checking against raster (non-SVG) content before treating this as settled.
            deck.ExportAsFixedFormat(os.path.abspath(pdf_path), 2)
            deck.Close()
            powerpoint.Quit()
        finally:
            pythoncom.CoUninitialize()
        return ok_result(row, f"PDF saved: {pdf_path}")
    except ImportError:
        return err_result(row, "save_pdf: comtypes not available — PDF export requires Windows + PowerPoint.")
    except Exception as e:
        return err_result(row, f"save_pdf: COM export failed: {e}")


# ---------------------------------------------------------------------------
# Dispatch map  —  function name -> callable
# ---------------------------------------------------------------------------

FUNCTION_MAP = {
    "create_ppt":               create_ppt,
    "set_default_populations":  set_default_populations,
    "insert_chart":             insert_chart,
    "insert_table":             insert_table,
    "insert_picture":           insert_picture,
    "insert_from_excel":        insert_from_excel,
    "open_excel":               open_excel,
    "close_excel":              close_excel,
    "update_text":              update_text,
    "empty_placeholder":        empty_placeholder,
    "save_ppt":                 save_ppt,
    "save_pdf":                 save_pdf,
}


# ---------------------------------------------------------------------------
# Run a complete Running Order
# ---------------------------------------------------------------------------

def run_running_order(rows: list[dict], settings: dict,
                      ctx: AssemblyContext = None) -> dict:
    """
    Execute a list of Running Order rows (already filtered to enabled only).

    settings dict must contain at minimum:
      ppt_template_path, cleaned_template_path, workfile_folder,
      reporting_unit_name, workfile_state

    Returns:
    {"status": "ok" | "error", "output_path": str, "elapsed": float, "log": list[dict]}
    """
    # Use a shared context if provided (batch run), otherwise create a fresh one.
    if ctx is None:
        ctx = AssemblyContext()

    workfile_state = settings.get("workfile_state")
    units = master_table_rows(workfile_state)
    ctx.report_context = build_report_context(settings, units)

    master_table_name = workfile_state.table_order[0] if workfile_state.table_order else ""
    reporting_row = None
    if ctx.report_context is not None:
        reporting_row = next(
            (r for r in units if str(r["unit_id"]) == ctx.report_context.unit_id), None
        )
    ctx.full_unit_set = (
        resolve_full_unit_set(reporting_row, master_table_name, workfile_state.tables)
        if reporting_row is not None else {}
    )

    t_start = time.perf_counter()

    normal_rows = [r for r in rows if str(r.get("scope", "normal")).strip() == "normal"]
    rows_to_run = normal_rows

    for i, row in enumerate(rows_to_run):
        func_name = str(row.get("function", "")).strip()

        func = FUNCTION_MAP.get(func_name)
        if func is None:
            result = err_result(row, f"Unknown function: '{func_name}'")
        else:
            try:
                result = func(ctx, row, settings)
            except Exception as e:
                result = err_result(row, f"Unhandled exception in '{func_name}': {traceback.format_exc()}")

        ctx.log.append(result)

        # Abort on error in structural functions
        if result["status"] == "error" and func_name in ("create_ppt",):
            ctx.log.append({"status": "aborted",
                            "message": "Batch aborted after create_ppt failure."})
            break

    elapsed = time.perf_counter() - t_start
    overall_status = "ok" if all(r["status"] in ("ok", "skip") for r in ctx.log) else "error"

    return {
        "status": overall_status,
        "output_path": ctx.output_path,
        "elapsed": elapsed,
        "log": ctx.log,
    }


# ---------------------------------------------------------------------------
# Private helpers — internal to insert_chart sub-steps
# ---------------------------------------------------------------------------

def _load_chart_data(cache_file: str, workfile_state=None):
    """Load a canonical data shape from the cache. Sub-step of insert_chart."""
    return load_shape(cache_file, workfile_state)


def _render_chart_image(base_chart_name: str, population_layers: list, width_emu: int, height_emu: int,
                        tweaks="", custom_chart_code=None):
    """
    Render a Matplotlib chart to SVG bytes sized to the placeholder.
    Sub-step of insert_chart. Returns image_bytes only — a Base Chart's
    only job. Statistics/unit lists are read directly off population_layers
    (already in scope here) by whatever needs them, e.g. Autotables
    (Feature List: Not built), rather than being relayed through render_chart.

    tweaks is the row's own tweaks column, passed straight through to the
    Base Chart function's tweaks parameter, uninterpreted here.

    base_chart_name is resolved built-in first, then against this workfile's
    own saved custom charts (get_chart_callable) — a custom chart behaves
    identically to a built-in from this point on. No report_context or any
    other runtime object is passed to a Base Chart function (Architecture,
    chart_inputs contract) — Selected-unit identity is already carried on
    the "Selected"-labelled entry in population_layers by the time this is
    called.

    width_emu/height_emu are passed straight through — EMU is the one real
    unit of size in this system (Architecture, Structural Design
    Principles); a Base Chart function converts to inches internally
    (divide by 914400) for its own matplotlib figsize. No percent
    conversion happens at this boundary any more.
    """
    chart_func = get_chart_callable(base_chart_name, custom_chart_code)
    return chart_func(population_layers, width_emu=width_emu, height_emu=height_emu, tweaks=tweaks)


def _insert_image_at_position(prs: Presentation, slide_index: int,
                               image_bytes, left_emu: int, top_emu: int,
                               width_emu: int, height_emu: int):
    """
    Insert an SVG image at the exact EMU position on the given slide, via
    the shared add_svg_picture dual-blip mechanism (see svg_insert.py) --
    every Base Chart returns SVG bytes (Architecture, SVG rendering
    methodology). Sub-step of insert_chart. Returns the inserted shape --
    insert_chart names it (CG_Chart_{row_id}) once it's back in scope
    there, rather than this function knowing about Running Order row
    identity at all.
    """
    if slide_index >= len(prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range "
            f"(template has {len(prs.slides)} slides)."
        )
    slide = prs.slides[slide_index]
    return add_svg_picture(
        slide, image_bytes.read(),
        left_emu, top_emu, width_emu, height_emu,
    )


# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------

def _ensure_output_folder(settings: dict) -> str:
    output_folder = settings.get("outputs_folder", "").strip()
    if not output_folder:
        workfile_folder = settings.get("workfile_folder", "").strip()
        output_folder = os.path.join(workfile_folder, "outputs") if workfile_folder else "outputs"
    os.makedirs(output_folder, exist_ok=True)
    return output_folder


def _safe_filename(name: str) -> str:
    """Strip characters that are unsafe in filenames."""
    import re
    safe = re.sub(r'[\\/:*?"<>|]', "_", name)
    return safe.strip("_") or "output"


def _int_or_none(value):
    try:
        return int(value)
    except (TypeError, ValueError):
        return None
