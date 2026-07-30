"""
insert_table.py
Running Order function: insert_table -- renders an Output Table (a grid of
constant text and resolved Stat Tag values, composited to a single image by
a Base Table function) and inserts it at the row's position. The Output
Table equivalent of insert_chart (assembly_engine.py), kept in its own
module under execution/tables/ rather than folded into assembly_engine,
mirroring how update_text was promoted to its own module under
execution/text/ (Architecture Decision 20).

An Output Table is defined independently of any Running Order row (the
same relationship a Stat Tag has to insert_chart rows) -- table_id anchors
an insert_table row to one of WorkfileState.output_tables, not the other
way round.

table_type_ref is resolved built-in first, then against this workfile's
own saved Custom Tables (get_table_callable) -- a custom table behaves
identically to a built-in from this point on, mirroring insert_chart's own
get_chart_callable resolution exactly.
"""

from pptx.util import Emu

from core.output_generation.execution.results import ok_result, err_result
from core.output_generation.execution.tables.resolve import resolve_output_table
from core.output_generation.execution.tables.custom_tables.resolve import get_table_callable


def insert_table(ctx, row: dict, settings: dict) -> dict:
    """Render an Output Table and insert it at the row's stored position."""
    if ctx.prs is None:
        return err_result(row, "insert_table: no open presentation (create_ppt not called?).")

    table_id = str(row.get("table_id", "") or "").strip()
    table_type_ref = str(row.get("table_type_ref", "") or "").strip()
    slide_index = _int_or_none(row.get("slide_index"))

    left_emu = _int_or_none(row.get("left_emu"))
    top_emu = _int_or_none(row.get("top_emu"))
    width_emu = _int_or_none(row.get("width_emu"))
    height_emu = _int_or_none(row.get("height_emu"))

    missing = []
    if not table_id: missing.append("table_id")
    if not table_type_ref: missing.append("table_type_ref")
    if slide_index is None: missing.append("slide_index")
    if None in (left_emu, top_emu, width_emu, height_emu):
        missing.append("position/size EMU values")
    if missing:
        return err_result(row, f"insert_table: missing required fields: {', '.join(missing)}")

    workfile_state = settings.get("workfile_state")
    grid_rows = workfile_state.output_tables.get(table_id) if workfile_state else None
    if not grid_rows:
        return err_result(row, f"insert_table: no Output Table found for table_id '{table_id}'.")

    resolved = resolve_output_table(grid_rows, workfile_state, ctx.full_unit_set or {})

    # Percent-of-the-7.5in-reference conversion a Base Table's width/height
    # parameters expect (matching the identical convention every Base Chart
    # uses) -- no ceiling or floor: a table's real placed size can
    # legitimately be smaller or larger than that reference, and clamping
    # the rendered resolution to it caused a genuine upscale (and visible
    # pixellation) whenever a table's real size exceeded it. See
    # Architecture, Structural Design Principles.
    NARROWER_EMU = 6858000
    width_pct = (width_emu / NARROWER_EMU) * 100
    height_pct = (height_emu / NARROWER_EMU) * 100

    tweaks = str(row.get("tweaks", "") or "").strip()
    custom_table_code = workfile_state.custom_table_code if workfile_state else {}
    try:
        table_func = get_table_callable(table_type_ref, custom_table_code)
        image_bytes = table_func(
            resolved["content"], resolved["column_widths"], resolved["row_heights"],
            width=width_pct, height=height_pct, tweaks=tweaks,
        )
    except Exception as e:
        return err_result(row, f"insert_table: render failed for '{table_type_ref}': {e}")

    try:
        slide = ctx.prs.slides[slide_index]
        image_bytes.seek(0)
        slide.shapes.add_picture(
            image_bytes, Emu(left_emu), Emu(top_emu), Emu(width_emu), Emu(height_emu),
        )
    except Exception as e:
        return err_result(row, f"insert_table: failed to insert image on slide {slide_index}: {e}")

    return ok_result(row, f"Output Table '{table_id}' inserted (slide {slide_index + 1})")


def _int_or_none(value):
    try:
        return int(value)
    except (TypeError, ValueError):
        return None
