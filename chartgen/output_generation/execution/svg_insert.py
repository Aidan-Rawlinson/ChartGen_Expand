"""
svg_insert.py

Inserts an SVG into a python-pptx presentation as a native SVG picture
rather than a raster one. The single mechanism for placing every Base Chart
and Base Table's output onto a slide, called from both assembly_engine.py
and tables/insert_table.py.

python-pptx's own add_picture() has no concept of SVG; it writes a single
raster blip relationship. PowerPoint stores an SVG picture as a dual-format
blip: a fallback raster, shown only by a viewer that does not understand the
extension, plus the SVG as a second part, wired together by an <a:extLst>
inside the picture's own <a:blip> pointing at the SVG part's relationship
id. python-pptx exposes no method for that extension, so it is hand-built
here against the lxml element tree.

Verified against python-pptx 1.0.2: a picture built this way round-trips
through Presentation.save()/Presentation(path) without error, and the SVG
part is correctly registered in [Content_Types].xml as an Override with
content type image/svg+xml. Confirmed working in real PowerPoint and both
PDF export pathways across the chart/table rollout.
"""

import io

from pptx.util import Emu
from pptx.opc.package import Part
from lxml import etree

# Fixed by the OOXML/Microsoft extension spec -- not something we choose.
_SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
_SVG_EXT_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
_SVG_CONTENT_TYPE = "image/svg+xml"
_SVG_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# A 1x1 transparent PNG used as the fallback raster when the caller
# doesn't supply one of its own. Only ever shown by a viewer that doesn't
# understand the SVG extension -- its own quality is irrelevant.
_BLANK_PNG_1PX = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
    "0000000a49444154789c6360000002000155e70da70000000049454e44ae426082"
)


def add_svg_picture(slide, svg_bytes: bytes, left: int, top: int, width: int, height: int,
                     fallback_png_bytes: bytes = None):
    """
    Insert svg_bytes as a native SVG picture on `slide`, at the given
    position/size (Emu ints -- same units and call shape as
    slide.shapes.add_picture, which this wraps for the raster half).

    Returns the picture shape, same object add_picture() would return,
    now also carrying the SVG extension pointing at the real image.
    """
    fallback = io.BytesIO(fallback_png_bytes or _BLANK_PNG_1PX)
    picture = slide.shapes.add_picture(fallback, Emu(left), Emu(top), Emu(width), Emu(height))

    slide_part = slide.part
    package = slide_part.package
    partname = package.next_partname("/ppt/media/svg%d.svg")
    svg_part = Part(partname, _SVG_CONTENT_TYPE, package, svg_bytes)
    svg_rel_id = slide_part.relate_to(svg_part, _SVG_REL_TYPE)

    blip = picture._element.blipFill.blip
    ext_lst = etree.SubElement(blip, f"{{{_A_NS}}}extLst")
    ext = etree.SubElement(ext_lst, f"{{{_A_NS}}}ext")
    ext.set("uri", _SVG_EXT_URI)
    svg_blip = etree.SubElement(ext, f"{{{_SVG_EXT_NS}}}svgBlip")
    svg_blip.set(f"{{{_R_NS}}}embed", svg_rel_id)

    return picture
